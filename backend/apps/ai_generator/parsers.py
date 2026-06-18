"""
Document text extractors for PDF and PowerPoint files.
Used as the first step in the AI content-generation pipeline.
"""
from __future__ import annotations

import re

from .exceptions import DocumentParseError

MAX_CHARS = 100_000
TRUNCATION_MARKER = "\n\n[CONTENIDO TRUNCADO — el documento supera el límite de procesamiento]"


def _clean(text: str) -> str:
    """Normalize whitespace and remove runs of blank lines."""
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r" {2,}", " ", text)
    return text.strip()


def _truncate(text: str) -> str:
    if len(text) <= MAX_CHARS:
        return text
    return text[:MAX_CHARS] + TRUNCATION_MARKER


def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract plain text from a PDF file using pdfplumber.
    Raises DocumentParseError if the file is empty, corrupt, or yields no text.
    """
    try:
        import pdfplumber  # noqa: PLC0415
    except ImportError as exc:
        raise DocumentParseError("pdfplumber no está instalado.") from exc

    try:
        with pdfplumber.open(file_path) as pdf:
            if not pdf.pages:
                raise DocumentParseError("El archivo PDF está vacío (sin páginas).")
            pages: list[str] = []
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                pages.append(page_text)
            raw = "\n\n".join(p for p in pages if p.strip())
    except DocumentParseError:
        raise
    except Exception as exc:
        raise DocumentParseError(f"No se pudo leer el PDF: {exc}") from exc

    text = _clean(raw)
    if not text:
        raise DocumentParseError("El archivo PDF no contiene texto extraíble.")

    return _truncate(text)


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract plain text from a PowerPoint file (.pptx) using python-pptx.
    Slides are separated by blank lines; text within each slide preserves paragraph order.
    Raises DocumentParseError if the file is empty, corrupt, or yields no text.
    """
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError as exc:
        raise DocumentParseError("python-pptx no está instalado.") from exc

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        raise DocumentParseError(f"No se pudo leer el archivo PowerPoint: {exc}") from exc

    slides: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            slides.append("\n".join(parts))

    if not slides:
        raise DocumentParseError(
            "El archivo PowerPoint no contiene texto extraíble."
        )

    raw = "\n\n".join(slides)
    text = _clean(raw)
    return _truncate(text)
